import os
import subprocess
import sys
import csv
import io
import textwrap
import base64
from langchain_core.tools import tool
from .utils import INTERNAL_SKILLS_DIR, USER_SKILLS_DIR, get_available_skills_hint, get_skill_suggestions, MEMORY_FILE, ensure_memory_exists, PROJECT_ROOT

# 尝试导入可选依赖
try:
    import docx
    import pypdf
    import openpyxl
    from pptx import Presentation
    HAS_OFFICE_DEPS = True
except ImportError:
    HAS_OFFICE_DEPS = False

@tool
def run_shell(command: str):
    """执行 Shell 命令。"""
    cmd_stripped = command.strip()
    if cmd_stripped.startswith("python3 ") or cmd_stripped.startswith("python "):
        parts = cmd_stripped.split(" ", 1)
        if len(parts) > 1:
            command = f"{sys.executable} {parts[1]}"
    try:
        result = subprocess.run(command, shell=True, capture_output=True, text=True, timeout=60)
        output = result.stdout
        if result.stderr:
            output += f"\nSTDERR: {result.stderr}"
        if len(output) > 2000:
            output = output[:2000] + "...(truncated)"
        return output
    except Exception as e:
        return f"命令执行错误: {e}"

@tool
def manage_skill(skill_name: str, action: str = "activate"):
    """
    管理特殊技能的生命周期。
    
    Args:
        skill_name (str): 技能名称（如 'excel_master', 'web_scraper'）。
        action (str): 操作类型。
            - 'activate': (默认) 激活技能。加载其 System Prompt 指令，使其具备特定领域的专业能力。
            - 'deactivate': 卸载技能。释放上下文空间，避免指令冲突或 Token 浪费。
    """
    action = action.lower()
    normalized_name = skill_name.strip()
    
    # === Action: Activate ===
    if action == "activate":
        search_paths = [
            os.path.join(INTERNAL_SKILLS_DIR, normalized_name, "SKILL.md"),
            os.path.join(USER_SKILLS_DIR, normalized_name, "SKILL.md")
        ]
        target_file = None
        skill_base_dir = None
        
        for path in search_paths:
            if os.path.exists(path):
                target_file = path
                skill_base_dir = os.path.dirname(path)
                break
        
        if target_file and skill_base_dir:
            try:
                with open(target_file, "r", encoding="utf-8") as f:
                    content = f.read()
                injected_content = content.replace("{SKILL_DIR}", skill_base_dir)
                return f"SYSTEM_INJECTION: {injected_content}"
            except Exception as e:
                return f"读取技能文件错误: {e}"
        else:
            suggestions = get_skill_suggestions(normalized_name)
            hint = get_available_skills_hint()
            err_msg = f"错误: 本地未找到技能 '{skill_name}'."
            if suggestions: err_msg += f"你是不是要找: {', '.join(suggestions)}."
            if hint: err_msg += f"可用技能: {hint}"
            return err_msg

    # === Action: Deactivate ===
    elif action == "deactivate":
        return f"SKILL_DEACTIVATION: {normalized_name}"
    
    else:
        return f"错误: 不支持的操作类型 '{action}'。请使用 'activate' 或 'deactivate'。"

def _read_docx(file_path, outline_only=False):
    doc = docx.Document(file_path)
    lines_pool = []
    outline = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines_pool.append("")
            continue
        if para.style.name.startswith('Heading'):
            try:
                level = int(para.style.name.split()[-1])
                outline.append(f"Line {len(lines_pool) + 1}: {'  ' * (level - 1)}- {text}")
            except: pass
        if len(text) > 120:
            wrapped = textwrap.fill(text, width=120)
            lines_pool.extend(wrapped.splitlines())
        else:
            lines_pool.append(text)
    img_count = len(doc.inline_shapes)
    if img_count > 0: lines_pool.append(f"\n[系统提示: 该文档包含 {img_count} 张图片]")
    if outline_only:
        if not outline: return "--- 文档大纲 ---\n[未检测到标准标题样式]\n"
        return "--- 文档大纲 (结构化导航) ---\n" + "\n".join(outline)
    return "\n".join(lines_pool)

def _read_pdf(file_path):
    reader = pypdf.PdfReader(file_path)
    text = []
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        wrapped_lines = []
        for line in page_text.splitlines():
             wrapped_lines.append(textwrap.fill(line, width=120) if len(line) > 120 else line)
        img_count = len(page.images)
        img_placeholder = f"\n[IMAGE_PLACEHOLDER: Page {i+1} 包含 {img_count} 张图片]\n" if img_count > 0 else ""
        text.append(f"--- Page {i+1} ---\n" + "\n".join(wrapped_lines) + img_placeholder)
    return "\n".join(text)

def _read_excel(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    output = []
    for sheet_name in wb.sheetnames:
        output.append(f"--- Sheet: {sheet_name} ---")
        ws = wb[sheet_name]
        si = io.StringIO()
        writer = csv.writer(si)
        for row in ws.rows:
            writer.writerow([cell.value for cell in row])
        output.append(si.getvalue())
    return "\n".join(output)

def _read_pptx(file_path):
    prs = Presentation(file_path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes: slide_text.append(f"[备注]: {notes}")
        if slide_text:
            text_content.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
    return "\n".join(text_content)

@tool
def read_file(file_path: str, start_line: int = 1, end_line: int = -1, outline_only: bool = False):
    """
    全能文件读取器。支持读取纯文本 (.txt, .md, .py, .json) 以及办公文档 (.docx, .pdf, .xlsx, .pptx)。
    1. 自动解析：对于 Office/PDF/PPT 文档，工具会自动提取文本内容及图片占位符信息。
    2. 大纲模式：设置 outline_only=True 获取 Docx 目录及行号映射，实现精准跳转。
    3. 分页功能：支持使用 start_line 和 end_line 进行按行分页读取（行号从 1 开始）。
    """
    if not os.path.exists(file_path): return f"错误: 未找到文件 '{file_path}'."
    ext = os.path.splitext(file_path)[1].lower()
    try:
        content_full = ""
        if HAS_OFFICE_DEPS:
            if ext == ".docx":
                content_full = _read_docx(file_path, outline_only=outline_only)
                if outline_only: return content_full
            elif ext == ".pdf": content_full = _read_pdf(file_path)
            elif ext in [".xlsx", ".xls"]: content_full = _read_excel(file_path)
            elif ext == ".pptx": content_full = _read_pptx(file_path)
        
        if not content_full:
            with open(file_path, 'r', encoding='utf-8') as f: lines = f.readlines()
        else:
            lines = content_full.splitlines(keepends=True)
            
        total_lines = len(lines)
        start_idx = max(0, start_line - 1)
        end_idx = min(end_line, total_lines) if end_line != -1 else min(start_idx + 500, total_lines)
        
        selected_lines = lines[start_idx:end_idx]
        header = f"--- 文件元数据 ---\n路径: {file_path}\n行数: {total_lines} | 当前范围: {start_idx+1}-{end_idx}\n"
        footer = f"\n\n[SYSTEM WARNING]: 文件未读完，后文还有 {total_lines - end_idx} 行。请调用 read_file(..., start_line={end_idx+1})." if end_idx < total_lines else ""
        return header + "\n" + "".join(selected_lines) + footer
    except Exception as e: return f"读取文件出错: {e}"

@tool
def write_file(file_path: str, content: str):
    """将文本内容写入指定文件（完全覆盖）。"""
    try:
        parent_dir = os.path.dirname(file_path)
        if parent_dir and not os.path.exists(parent_dir): os.makedirs(parent_dir, exist_ok=True)
        with open(file_path, 'w', encoding='utf-8') as f: f.write(content)
        return f"成功写入文件: {file_path}"
    except Exception as e: return f"写入文件出错: {e}"

@tool
def replace_in_file(file_path: str, old_string: str, new_string: str):
    """精确替换文件中的字符串。old_string 必须在文件中唯一存在。"""
    if not os.path.exists(file_path): return f"错误: 未找到文件 '{file_path}'."
    try:
        with open(file_path, 'r', encoding='utf-8') as f: content = f.read()
        if old_string not in content: return "错误: 在文件中未找到 old_string."
        if content.count(old_string) > 1: return "错误: old_string 不唯一."
        new_content = content.replace(old_string, new_string)
        with open(file_path, 'w', encoding='utf-8') as f: f.write(new_content)
        return f"成功在 {file_path} 中完成替换。"
    except Exception as e: return f"替换出错: {e}"

@tool
def save_memory(content: str):
    """
    保存长期记忆。用于记录用户偏好、重要事实或长期任务状态。
    会自动进行相似度检查，避免重复记录。
    """
    try:
        ensure_memory_exists()
        import datetime
        import difflib
        
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        # 智能去重检查
        for line in lines:
            content_part = line[20:].strip() if len(line) > 20 else line.strip()
            if difflib.SequenceMatcher(None, content_part, content.strip()).ratio() > 0.85:
                return f"记忆已存在 (相似度高)，跳过写入: {content}"
        
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
        entry = f"\n- [{timestamp}] {content}"
        with open(MEMORY_FILE, "a", encoding="utf-8") as f:
            f.write(entry)
        return f"成功保存记忆: {content}"
    except Exception as e: return f"保存记忆失败: {e}"

@tool
def forget_memory(content: str):
    """
    遗忘长期记忆。物理删除 MEMORY.md 中包含指定关键词的内容。
    用于修正错误或删除过时的用户偏好。
    """
    try:
        ensure_memory_exists()
        
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        new_lines = [l for l in lines if content not in l or l.startswith('#') or l.startswith('##')]
        deleted_count = len(lines) - len(new_lines)
        
        if deleted_count == 0:
            return f"未找到包含 '{content}' 的相关记忆。"
            
        with open(MEMORY_FILE, "w", encoding="utf-8") as f:
            f.writelines(new_lines)
        return f"成功遗忘 {deleted_count} 条相关记忆。"
    except Exception as e: return f"遗忘记忆失败: {e}"

@tool
def search_file(file_path: str, pattern: str, case_sensitive: bool = False):
    """
    在文件中搜索指定关键词或正则模式。支持 Office/PDF/PPT 文档。
    """
    if not os.path.exists(file_path): return f"错误: 未找到文件 '{file_path}'."
    ext = os.path.splitext(file_path)[1].lower()
    try:
        content_full = ""
        if HAS_OFFICE_DEPS:
            if ext == ".docx": content_full = _read_docx(file_path)
            elif ext == ".pdf": content_full = _read_pdf(file_path)
            elif ext in [".xlsx", ".xls"]: content_full = _read_excel(file_path)
            elif ext == ".pptx": content_full = _read_pptx(file_path)
        
        lines = content_full.splitlines() if content_full else open(file_path, 'r', encoding='utf-8').readlines()
        matches = []
        import re
        flags = 0 if case_sensitive else re.IGNORECASE
        for i, line in enumerate(lines):
            if re.search(pattern, line, flags):
                matches.append(f"Line {i+1}: {line.strip()}")
        if not matches: return f"未找到匹配 '{pattern}' 的内容."
        count = len(matches)
        result = f"--- 搜索结果 (共 {count} 处匹配) ---\n" + "\n".join(matches[:20])
        if count > 20: result += f"\n... (还有 {count-20} 处匹配已隐藏)"
        return result
    except Exception as e: return f"搜索出错: {e}"

@tool
def search_knowledge(query: str, collection: str = "documents"):
    """
    核心知识检索工具。
    功能：从本地向量库中搜索相关信息。
    适用场景：
    1. 查询已入库的文档（如白皮书、报价单）。 Collection: "documents"
    2. 回忆过去的对话历史（情景记忆）。 Collection: "episodic_memory"
    """
    # 动态定位脚本
    script_path = os.path.join(INTERNAL_SKILLS_DIR, "knowledge_base/scripts/query.py")
    if not os.path.exists(script_path):
        return "错误: 知识库技能脚本未找到。"
        
    cmd = [sys.executable, script_path, query, collection]
    try:
        # 注入 PYTHONPATH 确保脚本能找到 agent_core
        env = os.environ.copy()
        env["PYTHONPATH"] = PROJECT_ROOT
        
        res = subprocess.run(cmd, capture_output=True, text=True, env=env)
        if res.returncode != 0:
            return f"检索失败: {res.stderr}"
        return res.stdout
    except Exception as e:
        return f"执行错误: {e}"

@tool
def describe_image(image_path: str, prompt: str = "请详细描述这张图片的内容。"):
    """
    查看并分析本地图像文件（支持 PNG, JPG, WEBP）。
    注意：此工具会独立调用视觉模型 (gpt-4o-mini) 进行分析，不占用主模型的上下文。
    """
    if not os.path.exists(image_path):
        return f"错误: 未找到图片文件 '{image_path}'."
    
    try:
        # 获取图片后缀
        ext = os.path.splitext(image_path)[1].lower().strip('.')
        if ext == 'jpg': ext = 'jpeg'
        if ext not in ['png', 'jpeg', 'webp']:
            return f"错误: 不支持的图片格式 '{ext}'。请使用 png, jpg 或 webp。"

        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        
        # --- 视觉模型独立通道 ---
        # 支持任意兼容 OpenAI 接口的多模态模型 (如 Claude, Gemini, Qwen-VL 等)
        # 优先读取 VISION_LLM_X 系列变量
        from langchain_openai import ChatOpenAI
        from langchain_core.messages import HumanMessage
        
        vision_api_key = os.environ.get("VISION_LLM_API_KEY") or os.environ.get("OPENAI_API_KEY")
        vision_base_url = os.environ.get("VISION_LLM_BASE_URL")
        vision_model = os.environ.get("VISION_LLM_MODEL_NAME") or "gpt-4o-mini"
        
        if not vision_api_key:
            return "错误: 视觉能力未配置。请设置 VISION_LLM_API_KEY 环境变量 (支持 gpt-4o, claude-3-5, qwen-vl 等兼容 OpenAI 接口的模型)。"
            
        vision_llm = ChatOpenAI(
            model=vision_model,
            api_key=vision_api_key,
            temperature=0,
            base_url=vision_base_url # 如果为 None，ChatOpenAI 会默认使用官方地址
        )
        
        message = HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{ext};base64,{encoded_string}"},
                },
            ]
        )
        
        response = vision_llm.invoke([message])
        return f"--- 图像分析结果 ({image_path}) ---\n[Vision Model: {vision_model}]\n{response.content}"
    except Exception as e:
        return f"图像处理出错: {e}"

available_tools = [run_shell, manage_skill, read_file, write_file, replace_in_file, search_file, save_memory, forget_memory, search_knowledge, describe_image]
