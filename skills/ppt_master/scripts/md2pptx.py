import argparse
import os
import re
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def parse_markdown(md_content):
    """解析 Markdown 分页及内容"""
    raw_slides = re.split(r'^---+$', md_content, flags=re.MULTILINE)
    parsed_slides = []
    for raw in raw_slides:
        if not raw.strip(): continue
        slide_data = {
            'title': "Untitled", 
            'content': [], 
            'notes': "", 
            'visual_suggestions': [] # 新增：图示建议
        }
        lines = raw.strip().split('\n')
        mode = 'content'
        for line in lines:
            line = line.strip()
            if not line: continue
            
            if line.startswith('#'):
                slide_data['title'] = re.sub(r'^#+\s*(Slide\s*\d+[：|:|｜])?', '', line).strip()
            
            # 捕获图示建议
            elif "图示建议" in line or "建议画" in line:
                suggestion = line.replace("**图示建议：**", "").replace("- ", "").strip()
                slide_data['visual_suggestions'].append(suggestion)
                
            elif "Speaker Notes" in line or "演讲备注" in line:
                mode = 'notes'
            elif mode == 'content':
                # 过滤掉标签行
                if not line.startswith('**'):
                    slide_data['content'].append(line.lstrip('- ').lstrip('* '))
            elif mode == 'notes':
                slide_data['notes'] += line.lstrip('- ').lstrip('* ') + "\n"
        parsed_slides.append(slide_data)
    return parsed_slides

def find_layout_by_name(prs, name_keywords):
    for layout in prs.slide_layouts:
        for kw in name_keywords:
            if kw in layout.name:
                return layout
    return prs.slide_layouts[1] 

def clear_existing_slides(prs):
    """暴力清空模板自带的所有页面"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)

def add_visual_placeholder(slide, suggestions):
    """在页面右下角添加一个醒目的图示建议占位符"""
    if not suggestions: return
    
    # 创建一个圆角矩形
    left = Inches(5.5)
    top = Inches(2.5)
    width = Inches(4.0)
    height = Inches(3.0)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        left, top, width, height
    )
    # 样式：浅灰色填充，虚线边框
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    line = shape.line
    line.color.rgb = RGBColor(100, 100, 100)
    line.width = Pt(1.5)
    # line.dash_style = 
    
    # 文字
    tf = shape.text_frame
    tf.text = "🎨 图示建议区域\n\n" + "\n".join(suggestions)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.font.size = Pt(14)

def fill_slide(slide, slide_data):
    # 1. 标题
    if slide.shapes.title:
        slide.shapes.title.text = slide_data['title']
    
    # 2. 正文
    # 寻找正文框
    body_placeholders = [
        sp for shape in slide.placeholders 
        if (sp := shape).placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]
    ]
    body_placeholders.sort(key=lambda x: x.placeholder_format.idx)
    
    if body_placeholders:
        target = body_placeholders[0]
        tf = target.text_frame
        tf.clear()
        
        # 如果有图示建议，我们把正文框缩窄一点？(暂时不做，先只填文字)
        for point in slide_data['content']:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18) # 强制调整字号，避免太小
            p.space_after = Pt(10)
            
    # 3. 备注
    if slide_data['notes']:
        slide.notes_slide.notes_text_frame.text = slide_data['notes']
        
    # 4. 图示占位符
    add_visual_placeholder(slide, slide_data['visual_suggestions'])

def create_ppt(slides, output_path, template_path):
    if not os.path.exists(template_path):
        print(f"⚠️ 未找到模板: {template_path}，已回退为空白模板。")
        prs = Presentation()
    else:
        prs = Presentation(template_path)
        # clear_existing_slides(prs) # 暂时禁用清空逻辑，防止破坏企业模板的底层结构导致文件损坏

    for i, slide_data in enumerate(slides):
        if i == 0:
            layout = find_layout_by_name(prs, ["封面"])
        elif any(kw in slide_data['title'] for kw in ["目录", "提纲"]):
            layout = find_layout_by_name(prs, ["目录", "提纲"])
        elif any(kw in slide_data['title'] for kw in ["总结", "谢", "Thanks"]):
            layout = find_layout_by_name(prs, ["封底", "结束"])
        else:
            layout = find_layout_by_name(prs, ["内页", "正文"])
            
        slide = prs.slides.add_slide(layout)
        try:
            fill_slide(slide, slide_data)
        except Exception as e:
            print(f"Warning: Failed to fill slide {i}: {e}")

    prs.save(output_path)
    print(f"🎉 成功生成: {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input")
    parser.add_argument("output")
    parser.add_argument("--template", default="skills/ppt_master/templates/2024金蝶集团PPT模板.pptx")
    args = parser.parse_args()
    
    with open(args.input, 'r', encoding='utf-8') as f:
        md = f.read()
    create_ppt(parse_markdown(md), args.output, args.template)
